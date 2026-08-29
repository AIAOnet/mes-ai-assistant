"""Persistent, permission-aware document ingestion and lexical retrieval for MES RAG."""
from __future__ import annotations
import hashlib, html, io, json, math, os, re, threading, urllib.error, urllib.request, zipfile
from datetime import datetime, timezone
from email import policy
from email.parser import BytesParser
from pathlib import Path
from xml.etree import ElementTree

MAX_UPLOAD_BYTES = 25 * 1024 * 1024
TEXT_EXTENSIONS = {".txt",".md",".markdown",".rst",".log",".csv",".tsv",".json",".jsonl",".yaml",".yml",".toml",".ini",".cfg",".conf",".xml",".html",".htm",".css",".js",".jsx",".ts",".tsx",".py",".java",".c",".h",".cpp",".hpp",".cs",".go",".rs",".sql",".sh",".ps1",".bat",".properties"}
SUPPORTED_EXTENSIONS = sorted(TEXT_EXTENSIONS | {".pdf",".docx",".xlsx",".xls",".pptx",".odt",".ods",".odp",".rtf",".eml",".epub",".zip"})
class KnowledgeValidationError(ValueError): pass

class EmbeddingError(RuntimeError): pass

class EmbeddingClient:
    """Small OpenAI/Opper-compatible embeddings client configured only by environment."""
    def __init__(self, endpoint="", api_key="", model="", auth_header="Authorization", timeout=45):
        self.endpoint, self.api_key, self.model = endpoint.strip(), api_key.strip(), model.strip()
        self.auth_header, self.timeout = auth_header.strip() or "Authorization", float(timeout)
    @classmethod
    def from_environment(cls):
        return cls(os.getenv("MES_RAG_EMBEDDING_ENDPOINT", ""),
                   os.getenv("MES_RAG_EMBEDDING_API_KEY", "") or os.getenv("MES_AI_API_KEY", ""),
                   os.getenv("MES_RAG_EMBEDDING_MODEL", ""),
                   os.getenv("MES_RAG_EMBEDDING_AUTH_HEADER", "Authorization"),
                   os.getenv("MES_RAG_EMBEDDING_TIMEOUT_SECONDS", "45"))
    @property
    def configured(self): return bool(self.endpoint and self.api_key and self.model)
    def embed(self, texts):
        if not self.configured: raise EmbeddingError("Semantic embeddings are not configured")
        if not texts: return []
        vectors=[]
        for start in range(0,len(texts),32):
            payload=json.dumps({"model":self.model,"input":texts[start:start+32]}).encode()
            auth=self.api_key if self.auth_header.lower()!="authorization" else f"Bearer {self.api_key}"
            request=urllib.request.Request(self.endpoint,data=payload,headers={"Content-Type":"application/json",self.auth_header:auth},method="POST")
            try:
                with urllib.request.urlopen(request,timeout=self.timeout) as response: body=json.loads(response.read())
                items=body.get("data",body.get("embeddings",[]))
                batch=[item.get("embedding",item) if isinstance(item,dict) else item for item in items]
                if len(batch)!=len(texts[start:start+32]) or any(not isinstance(x,list) for x in batch): raise EmbeddingError("Embedding service returned an invalid response")
                vectors.extend(batch)
            except (OSError,ValueError,urllib.error.HTTPError) as error: raise EmbeddingError(f"Embedding service unavailable: {error}") from error
        return vectors

def _clean(text):
    text = html.unescape(text).replace("\x00", " "); text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"[ \t]+", " ", text); return re.sub(r"\n{3,}", "\n\n", text).strip()[:6_000_000]

def _archive_text(data):
    results, total = [], 0
    with zipfile.ZipFile(io.BytesIO(data)) as archive:
        members = [x for x in archive.infolist() if not x.is_dir()]
        if len(members) > 250: raise KnowledgeValidationError("Archive contains too many files")
        for member in members:
            total += member.file_size
            if total > MAX_UPLOAD_BYTES * 4: raise KnowledgeValidationError("Archive expands beyond the safe extraction limit")
            suffix = Path(member.filename).suffix.lower()
            if suffix in TEXT_EXTENSIONS or suffix == ".xhtml": results.append(f"[{member.filename}]\n{_clean(archive.read(member).decode('utf-8', errors='replace'))}")
            elif member.filename.endswith("content.xml"):
                root = ElementTree.fromstring(archive.read(member)); results.append("\n".join(x.strip() for x in root.itertext() if x.strip()))
    return _clean("\n\n".join(results))

def extract_text(filename, data):
    if not data: raise KnowledgeValidationError("The uploaded file is empty")
    if len(data) > MAX_UPLOAD_BYTES: raise KnowledgeValidationError("File exceeds the 25 MB upload limit")
    suffix = Path(filename).suffix.lower()
    try:
        if suffix in TEXT_EXTENSIONS: return _clean(data.decode("utf-8", errors="replace"))
        if suffix == ".pdf":
            from pypdf import PdfReader
            return _clean("\n".join(x.extract_text() or "" for x in PdfReader(io.BytesIO(data)).pages))
        if suffix == ".docx":
            from docx import Document
            doc=Document(io.BytesIO(data)); parts=[x.text for x in doc.paragraphs]; parts += [" | ".join(c.text for c in row.cells) for table in doc.tables for row in table.rows]; return _clean("\n".join(parts))
        if suffix == ".xlsx":
            from openpyxl import load_workbook
            book=load_workbook(io.BytesIO(data),read_only=True,data_only=True); return _clean("\n".join(f"[{s.title}] "+" | ".join("" if v is None else str(v) for v in row) for s in book.worksheets for row in s.iter_rows(values_only=True)))
        if suffix == ".xls":
            import xlrd
            book=xlrd.open_workbook(file_contents=data); return _clean("\n".join(f"[{s.name}] "+" | ".join(str(s.cell_value(r,c)) for c in range(s.ncols)) for s in book.sheets() for r in range(s.nrows)))
        if suffix == ".pptx":
            from pptx import Presentation
            return _clean("\n".join(x.text for slide in Presentation(io.BytesIO(data)).slides for x in slide.shapes if hasattr(x,"text")))
        if suffix == ".eml":
            msg=BytesParser(policy=policy.default).parsebytes(data); parts=[f"Subject: {msg.get('subject','')}",f"From: {msg.get('from','')}"]+[x.get_content() for x in msg.walk() if x.get_content_type()=="text/plain"]; return _clean("\n".join(parts))
        if suffix == ".rtf": return _clean(re.sub(r"\\[a-zA-Z]+-?\d* ?|[{}]", " ", data.decode("latin-1",errors="replace")))
        if suffix in {".odt",".ods",".odp",".epub",".zip"}: return _archive_text(data)
    except KnowledgeValidationError: raise
    except Exception as error: raise KnowledgeValidationError(f"Could not extract text from {suffix or 'this file'}: {error}") from error
    raise KnowledgeValidationError(f"Unsupported format '{suffix or 'unknown'}'. Supported: {', '.join(SUPPORTED_EXTENSIONS)}")

def _terms(text): return re.findall(r"[a-z0-9][a-z0-9_-]{1,}", text.lower())
def _chunks(text,size=1400,overlap=180):
    chunks,current=[],""
    for paragraph in [x.strip() for x in re.split(r"\n\s*\n",text) if x.strip()]:
        while len(paragraph)>size: chunks.append(paragraph[:size]); paragraph=paragraph[size-overlap:]
        if current and len(current)+len(paragraph)+2>size: chunks.append(current); current=current[-overlap:]+"\n\n"+paragraph
        else: current=f"{current}\n\n{paragraph}".strip()
    if current: chunks.append(current)
    return chunks[:1500]

class KnowledgeStore:
    def __init__(self,root="rag_data",embedder=None):
        self.root=Path(root); self.files=self.root/"documents"; self.index_path=self.root/"index.json"; self._lock=threading.RLock(); self.files.mkdir(parents=True,exist_ok=True)
        self.embedder=embedder or EmbeddingClient.from_environment(); self.last_embedding_error=""
        self.search_mode=os.getenv("MES_RAG_SEARCH_MODE","hybrid").strip().lower()
        try: self.semantic_weight=max(0.0,min(1.0,float(os.getenv("MES_RAG_SEMANTIC_WEIGHT","0.65"))))
        except ValueError: self.semantic_weight=0.65
        try: self._index=json.loads(self.index_path.read_text(encoding="utf-8"))
        except (OSError,json.JSONDecodeError): self._index={"documents":[]}
    @staticmethod
    def _public(x): return {k:v for k,v in x.items() if k not in {"chunks","embeddings","stored_name"}}|{"chunk_count":len(x.get("chunks",[])),"semantic_indexed":len(x.get("embeddings",[]))==len(x.get("chunks",[])) and bool(x.get("chunks"))}
    def _save(self):
        temp=self.index_path.with_suffix(".tmp"); temp.write_text(json.dumps(self._index,ensure_ascii=False,indent=2),encoding="utf-8"); temp.replace(self.index_path)
    def add(self,filename,data,title="",version="",machine_id="",alarm_type="",roles=None):
        safe=Path(filename).name[:180]; text=extract_text(safe,data)
        if len(_terms(text))<3: raise KnowledgeValidationError("No useful searchable text was found in the file")
        allowed=sorted(set(roles or ["admin","operator","viewer"])&{"admin","operator","viewer"})
        if not allowed: raise KnowledgeValidationError("At least one valid permission role is required")
        digest=hashlib.sha256(data).hexdigest(); doc_id=digest[:20]
        with self._lock:
            if any(x["id"]==doc_id for x in self._index["documents"]): raise KnowledgeValidationError("This exact file is already indexed")
            stored=f"{doc_id}{Path(safe).suffix.lower()}"; (self.files/stored).write_bytes(data)
            chunks=_chunks(text); embeddings=[]
            if self.embedder.configured and self.search_mode=="hybrid":
                try: embeddings=self.embedder.embed(chunks); self.last_embedding_error=""
                except EmbeddingError as error: self.last_embedding_error=str(error)
            record={"id":doc_id,"filename":safe,"stored_name":stored,"title":title.strip()[:200] or safe,"version":version.strip()[:80],"machine_id":machine_id.strip().upper()[:40],"alarm_type":alarm_type.strip()[:100],"roles":allowed,"size_bytes":len(data),"sha256":digest,"uploaded_at":datetime.now(timezone.utc).isoformat(),"chunks":chunks,"embeddings":embeddings}
            self._index["documents"].append(record); self._save()
        return self._public(record)
    def list(self,role="admin"):
        with self._lock: return [self._public(x) for x in reversed(self._index["documents"]) if role in x["roles"]]
    def get(self,doc_id,role="admin"):
        with self._lock:
            record=next((x for x in self._index["documents"] if x["id"]==doc_id and role in x["roles"]),None)
            return self._public(record) if record else None
    def delete(self,doc_id):
        with self._lock:
            record=next((x for x in self._index["documents"] if x["id"]==doc_id),None)
            if not record:return False
            self._index["documents"].remove(record); (self.files/record["stored_name"]).unlink(missing_ok=True); self._save(); return True
    @staticmethod
    def _cosine(a,b):
        if not a or not b or len(a)!=len(b): return 0.0
        denominator=math.sqrt(sum(x*x for x in a))*math.sqrt(sum(x*x for x in b))
        return sum(x*y for x,y in zip(a,b))/denominator if denominator else 0.0
    def status(self):
        with self._lock:
            documents=len(self._index["documents"]); embedded=sum(1 for x in self._index["documents"] if len(x.get("embeddings",[]))==len(x.get("chunks",[])) and x.get("chunks"))
        mode="hybrid" if self.embedder.configured and self.search_mode=="hybrid" else "keyword"
        return {"mode":mode,"embedding_configured":self.embedder.configured,"semantic_weight":self.semantic_weight,"embedding_model":self.embedder.model or None,"documents":documents,"embedded_documents":embedded,"pending_documents":documents-embedded,"last_embedding_error":self.last_embedding_error or None}
    def reindex(self):
        if not self.embedder.configured: raise EmbeddingError("Semantic embeddings are not configured")
        indexed=0
        with self._lock:
            for document in self._index["documents"]:
                if len(document.get("embeddings",[]))==len(document.get("chunks",[])) and document.get("chunks"): continue
                document["embeddings"]=self.embedder.embed(document["chunks"]); indexed+=1
            self._save(); self.last_embedding_error=""
        return indexed
    def search(self,query,role,limit=5,machine_id="",alarm_type=""):
        terms=_terms(query); matches=[]; query_vector=[]
        if not terms:return []
        if self.embedder.configured and self.search_mode=="hybrid":
            try: query_vector=self.embedder.embed([query])[0]; self.last_embedding_error=""
            except EmbeddingError as error: self.last_embedding_error=str(error)
        with self._lock: documents=list(self._index["documents"])
        for doc in documents:
            if role not in doc["roles"]:continue
            if machine_id and doc["machine_id"] and doc["machine_id"]!=machine_id.upper():continue
            if alarm_type and doc["alarm_type"] and doc["alarm_type"].lower()!=alarm_type.lower():continue
            vectors=doc.get("embeddings",[])
            for index,chunk in enumerate(doc["chunks"]):
                tokens=_terms(chunk); lexical=sum(1+math.log(tokens.count(t)) for t in set(terms) if t in tokens)
                if query.lower().strip() in chunk.lower():lexical+=4
                semantic=self._cosine(query_vector,vectors[index]) if query_vector and index<len(vectors) else 0.0
                if lexical or semantic>0:
                    matches.append({"lexical_score":lexical,"semantic_score":semantic,"text":chunk,"chunk":index,"document":self._public(doc)})
        max_lexical=max((x["lexical_score"] for x in matches),default=1) or 1
        hybrid=bool(query_vector)
        for item in matches:
            lexical=item.pop("lexical_score")/max_lexical; semantic=max(0.0,item["semantic_score"])
            item["score"]=round(((1-self.semantic_weight)*lexical+self.semantic_weight*semantic) if hybrid else lexical,4)
            item["semantic_score"]=round(semantic,4); item["search_mode"]="hybrid" if hybrid else "keyword"
        return sorted(matches,key=lambda x:x["score"],reverse=True)[:max(1,min(int(limit),10))]
